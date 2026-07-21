"""文档在线预览 — 上传 pdf/docx/xlsx/pptx，直接在线渲染预览。"""
from __future__ import annotations

import base64 as _b64
import io
import logging
from pathlib import Path

from flask import Blueprint, current_app, jsonify, render_template, request, send_file
from pypdf import PdfReader, PdfWriter

from auth.decorators import commit_usage, remaining_for, require_usage
from extensions import limiter
from utils.helpers import safe_download_path, safe_filename

logger = logging.getLogger(__name__)
tool_bp = Blueprint("doc_viewer", __name__)

_ALLOWED_EXTS = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


@tool_bp.get("/")
def index():
    return render_template(
        "tools_base.html",
        tool={"id": "doc_viewer", "name": "文档在线预览", "icon": "bi-eye", "color": "#20c997"},
        remaining=remaining_for("doc_viewer"),
        body_template="tools/doc_viewer/_body.html",
    )


@tool_bp.post("/process")
@limiter.limit(lambda: current_app.config["RATELIMIT_TOOL"])
@require_usage("doc_viewer")
def process():
    file = request.files.get("file")
    if file is None or not file.filename:
        commit_usage("doc_viewer", success=False, message="未选择文件")
        return jsonify(error="请选择文件"), 400

    ext = Path(file.filename).suffix.lower().lstrip(".")
    if ext not in _ALLOWED_EXTS:
        return jsonify(error="仅支持 PDF / DOCX / XLSX / PPTX 格式"), 400

    raw = file.read()
    if len(raw) > 50 * 1024 * 1024:  # 50MB limit
        return jsonify(error="文件太大（最大 50MB）"), 400

    try:
        if ext == "pdf":
            result = _preview_pdf(raw, file.filename)
        elif ext == "docx":
            result = _preview_docx(raw)
        elif ext == "xlsx":
            result = _preview_xlsx(raw)
        elif ext == "pptx":
            result = _preview_pptx(raw)
        else:
            return jsonify(error=f"暂不支持 {ext} 格式"), 400

        commit_usage("doc_viewer", success=True)
        return jsonify(ok=True, format=ext, **result)
    except Exception as exc:
        logger.exception("doc_viewer failed for %s: %s", ext, exc)
        commit_usage("doc_viewer", success=False, message=str(exc))
        return jsonify(error=f"预览失败：{exc}"), 500


@tool_bp.get("/download/<filename>")
def download(filename: str):
    upload_dir: Path = current_app.config["UPLOAD_DIR"]
    target = safe_download_path(upload_dir, filename)
    if target is None or not target.exists():
        return jsonify(error="文件不存在"), 404
    return send_file(target, as_attachment=True)


# ── converters ──────────────────────────────────────────────────────────

def _preview_pdf(raw: bytes, filename: str) -> dict:
    """使用 PyMuPDF 将每页渲染为 base64 PNG。"""
    import fitz  # PyMuPDF  # noqa: PLC0415

    doc = fitz.open(stream=raw, filetype="pdf")
    try:
        pages = []
        for i in range(doc.page_count):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=150)
            img_b64 = _b64.b64encode(pix.tobytes("png")).decode("ascii")
            pages.append({
                "num": i + 1,
                "width": pix.width,
                "height": pix.height,
                "src": f"data:image/png;base64,{img_b64}",
            })
        return {"pages": pages, "page_count": len(pages), "filename": filename}
    finally:
        doc.close()


def _preview_docx(raw: bytes) -> dict:
    """把 docx 段落提取为 HTML。"""
    from docx import Document  # noqa: PLC0415

    doc = Document(io.BytesIO(raw))
    html_parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            html_parts.append("<br>")
            continue
        style = para.style.name if para.style else ""
        if style.startswith("Heading"):
            level = "".join(c for c in style if c.isdigit()) or "1"
            html_parts.append(f"<h{level}>{_esc(text)}</h{level}>")
        else:
            html_parts.append(f"<p>{_esc(text)}</p>")

    # tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = "".join(f"<td>{_esc(cell.text)}</td>" for cell in row.cells)
            rows.append(f"<tr>{cells}</tr>")
        html_parts.append(f"<table class='table table-bordered table-sm'>{''.join(rows)}</table>")

    return {"html": "".join(html_parts), "format_name": "Word (DOCX)"}


def _preview_xlsx(raw: bytes) -> dict:
    """把 xlsx 渲染为 HTML 表格（最多前 10 个 sheet，每个最多 100×20 格）。"""
    from openpyxl import load_workbook  # noqa: PLC0415

    wb = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    sheets_html = []
    for si, name in enumerate(wb.sheetnames):
        if si >= 10:
            break
        ws = wb[name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        rows = rows[:100]
        max_cols = max((len(r) for r in rows if r), default=0)
        html = '<table class="table table-bordered table-sm table-striped small mb-3">'
        for ri, row in enumerate(rows):
            html += "<tr>"
            for ci in range(min(max_cols, 20)):
                val = str(row[ci]) if ci < len(row) and row[ci] is not None else ""
                tag = "th" if ri == 0 else "td"
                html += f"<{tag}>{_esc(val)}</{tag}>"
            html += "</tr>"
        html += "</table>"
        sheets_html.append(f"<h6>Sheet: {_esc(name)}</h6>{html}")
    wb.close()
    return {"html": "".join(sheets_html), "format_name": "Excel (XLSX)"}


def _preview_pptx(raw: bytes) -> dict:
    """把 pptx 幻灯片文字提取为 HTML。"""
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(raw))
    slides_html = []
    for si, slide in enumerate(prs.slides):
        slide_html = f'<div class="card mb-2"><div class="card-header small fw-bold">幻灯片 {si + 1}</div><div class="card-body py-2">'
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(f"<p class='mb-1'>{_esc(t)}</p>")
        slide_html += "".join(texts) or "<p class='text-muted'>（无文字内容）</p>"
        slide_html += "</div></div>"
        slides_html.append(slide_html)
    return {"html": "".join(slides_html), "format_name": "PowerPoint (PPTX)"}


def _esc(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
